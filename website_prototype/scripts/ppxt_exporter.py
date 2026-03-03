import os
from pptx import Presentation
from pptx.util import Emu  # On utilise Emu pour une précision au pixel près
from PIL import Image
from pathlib import Path


def create_pptx_from_images(input_folder, output_path):
    # 1. Vérification de l'existence du dossier
    if not input_folder.exists() or not input_folder.is_dir():
        print(f"Erreur : Le dossier '{input_folder}' est introuvable.")
        return

    # 2. Lister et trier les fichiers PNG
    files = [f for f in os.listdir(input_folder) if f.lower().endswith('.png')]
    files.sort()

    if not files:
        print("Aucune image PNG trouvée dans le dossier.")
        return

    # 3. Initialisation de la présentation
    prs = Presentation()

    try:
        # On utilise la première image pour définir la taille des diapositives
        first_img_path = input_folder / files[0]
        with Image.open(first_img_path) as img:
            width_px, height_px = img.size

        # Conversion des pixels en EMU (1 pixel standard = 9525 EMU)
        width_emu = Emu(width_px * 9525 // 72)  # Ajustement standard 72/96 dpi
        # Plus simple : python-pptx gère mieux si on définit directement les dimensions
        prs.slide_width = width_px * 9525
        prs.slide_height = height_px * 9525

        for file in files:
            img_path = input_folder / file

            # Layout 6 = Slide vide
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            # Insertion de l'image en forçant la taille d'origine
            # On passe le chemin en string pour la compatibilité avec python-pptx
            slide.shapes.add_picture(
                str(img_path),
                0, 0,
                width=prs.slide_width,
                height=prs.slide_height
            )

        # 4. Sauvegarde
        prs.save(str(output_path))
        print(f"\n--- Succès ! ---")
        print(f"Le PowerPoint a été généré ici : {output_path}")

    except Exception as e:
        print(f"Une erreur est survenue : {e}")


# --- Configuration interactive avec ta logique ---

dossier_input = input("Veuillez entrer le chemin du dossier contenant les images : ").strip().strip('"')
path_image = Path(dossier_input)

nom_sortie = input("Quel nom pour le nuancier PPTX à transmettre à l'équipe ? ").strip()

if not nom_sortie.lower().endswith(".pptx"):
    nom_sortie += ".pptx"

chemin_final = path_image / nom_sortie

create_pptx_from_images(path_image, chemin_final)